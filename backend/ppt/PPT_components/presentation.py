"""
PowerPoint presentation generation and manipulation functions.
"""

import math
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN

from .utils import analyze_format, extract_number, format_number

def remove_bullets_from_paragraph(paragraph):
    """Remove bullet formatting from a paragraph completely"""
    from lxml import etree
    
    # Set level to 0
    paragraph.level = 0
    
    # Get or create paragraph properties
    pPr = paragraph._element.get_or_add_pPr()
    
    # Namespace for DrawingML elements
    nsdecls = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
    
    # Remove all bullet-related elements
    for child in list(pPr):
        tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
        # Remove any bullet-related elements
        if tag in ['buChar', 'buFont', 'buAutoNum', 'buBlip', 'buNone', 'buStop']:
            pPr.remove(child)
    
    # Explicitly add buNone to disable bullets
    # First remove if it exists to avoid duplicates
    try:
        buNone_elem = pPr.find('a:buNone', nsdecls)
        if buNone_elem is not None:
            pPr.remove(buNone_elem)
    except:
        pass
    
    # Create and insert buNone element at the beginning
    buNone = etree.Element('{http://schemas.openxmlformats.org/drawingml/2006/main}buNone')
    pPr.insert(0, buNone)

def remove_bullets_from_text_frame(text_frame):
    """Remove bullet formatting from all paragraphs in a text frame"""
    for paragraph in text_frame.paragraphs:
        remove_bullets_from_paragraph(paragraph)

def format_cost_to_millions(cost_value, decimal_places=1, prefix="$", suffix="M"):
    """
    Convert a cost value to millions format.
    
    Args:
        cost_value (float): The cost value to convert
        decimal_places (int): Number of decimal places
        prefix (str): Currency prefix (e.g., "$")
        suffix (str): Suffix (e.g., "M" for millions)
        
    Returns:
        str: Formatted cost string (e.g., "$0.1M")
    """
    if cost_value == 0:
        return f"{prefix}0{suffix}"
    
    millions = cost_value / 1_000_000
    format_str = f"{{:.{decimal_places}f}}"
    formatted = format_str.format(millions)
    return f"{prefix}{formatted}{suffix}"

def add_slide_notes(slide, notes_text):
    """Add notes to a slide"""
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = notes_text

def calculate_text_width(text, font_size):
    """Calculate approximate width of text in points based on font size"""
    # This is an approximation. Actual width may vary based on font and characters
    return len(text) * font_size * 0.6  # 0.6 is an approximation factor for width/height ratio

def split_text_for_pill(text, available_width, font_size):
    """Split text to fit in available width, at a space for natural word breaks"""
    if calculate_text_width(text, font_size) <= available_width:
        return text, False  # Text fits, no need to split
        
    words = text.split()
    if len(words) == 1:
        return text, True  # Single word, can't split, needs shortening
        
    # Find all potential split points (spaces)
    split_points = []
    current_pos = 0
    for i, char in enumerate(text):
        if char == ' ':
            # Calculate parts if we split here
            first_part = text[:i].strip()
            second_part = text[i+1:].strip()
            width1 = calculate_text_width(first_part, font_size)
            width2 = calculate_text_width(second_part, font_size)
            
            # Only consider split if both parts fit within available width
            if width1 <= available_width and width2 <= available_width:
                # Calculate how balanced the split is (closer to 0 is more balanced)
                balance_score = abs(width1 - width2)
                # Calculate how close to middle the split is (closer to 0 is better)
                middle_score = abs(i - len(text)/2)
                # Combine scores (you can adjust weights if needed)
                total_score = balance_score + middle_score
                split_points.append((i, total_score))
    
    if not split_points:
        return text, True  # No good split points found, needs shortening
    
    # Find the split point with the lowest score
    best_split = min(split_points, key=lambda x: x[1])[0]
    
    # Split at best point
    first_part = text[:best_split].strip()
    second_part = text[best_split+1:].strip()
    
    # If best split still results in a part that's too wide, text needs shortening
    if (calculate_text_width(first_part, font_size) > available_width or
        calculate_text_width(second_part, font_size) > available_width):
        return text, True  # Text needs shortening
        
    # Return parts separately to be handled by the text frame
    return (first_part, second_part), False

def draw_group_boxes(prs, slide, df, config, shortened_text_map, slide_title=None):
    """Draw group boxes on a PowerPoint slide"""
    # Get slide dimensions
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    
    # Update the slide title placeholder if provided
    if slide_title:
        # Find and update the title shape (usually at index 0 in shapes)
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.has_text_frame:
                # Check if this is a title placeholder
                if "title" in shape.name.lower() or (hasattr(shape, "is_placeholder") and shape.is_placeholder):
                    text_frame = shape.text_frame
                    text_frame.clear()
                    p = text_frame.paragraphs[0]
                    p.text = slide_title
                    p.font.size = Pt(24)
                    p.font.bold = True
                    p.font.name = "Arial"
                    p.font.color.rgb = RGBColor(0, 0, 0)
                    p.alignment = PP_ALIGN.LEFT  # Left-aligned
                    remove_bullets_from_paragraph(p)
                    break
    
    # Draw the legend
    legend_left = Inches(0.5)  # Position from left
    legend_top = Inches(1.1)   # Position from top
    legend_height = Inches(0.25)  # Height of legend boxes
    legend_width = Inches(1.5)   # Width of legend boxes
    legend_spacing = Inches(0.2)  # Spacing between legend items
    
    # Add "Legend:" text
    legend_text = slide.shapes.add_textbox(
        legend_left, 
        legend_top, 
        Inches(0.6),  # Width for "Legend:" text
        legend_height
    )
    legend_text_frame = legend_text.text_frame
    p = legend_text_frame.paragraphs[0]
    p.text = "Legend:"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.name = "Arial"
    remove_bullets_from_paragraph(p)
    
    # Start position for legend boxes
    box_start_x = legend_left + Inches(0.7)  # After "Legend:" text
    
    # Get color ranges from config
    color_ranges = config.get("color_ranges", {})
    
    # Colors in order
    colors = ["Green", "Yellow", "Red"]
    
    # Calculate total count (sum of all L2 counts, ignoring '-')
    total_count = df.loc[df['Count'] != '-', 'Count'].astype(int).sum()
    
    # Add total count text
    total_count_text = slide.shapes.add_textbox(
        slide_width - Inches(1.87),  # Position 0.37" from right edge
        legend_top,   # At the same level as legend
        Inches(1.5),  # Width for count text
        Inches(0.3)   # Height for count text
    )
    count_text_frame = total_count_text.text_frame
    count_p = count_text_frame.paragraphs[0]
    count_p.text = f"N = {total_count}"
    count_p.font.size = Pt(12)
    count_p.font.bold = True
    count_p.font.name = "Arial"
    count_p.alignment = PP_ALIGN.RIGHT  # Right align the text
    remove_bullets_from_paragraph(count_p)
    
    # Draw legend boxes
    for i, color in enumerate(colors):
        box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            box_start_x + (i * (legend_width + legend_spacing)),
            legend_top,
            legend_width,
            legend_height
        )
        
        # Set box color using hex code from color_map
        hex_color = config["color_map"].get(color, "000000")
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(
            int(hex_color[0:2], 16),
            int(hex_color[2:4], 16),
            int(hex_color[4:6], 16)
        )
        
        # Add text for range
        min_val, max_val = color_ranges.get(color, (0, 0))
        if color == "Red":
            range_text = f"≥ {min_val} Applications"
        else:
            range_text = f"{min_val}-{max_val} Applications"
            
        text_frame = box.text_frame
        p = text_frame.paragraphs[0]
        p.text = range_text
        p.font.size = Pt(9)
        p.font.bold = True
        p.font.name = "Arial"
        p.alignment = PP_ALIGN.CENTER
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        remove_bullets_from_paragraph(p)
    
    # Extract configuration parameters
    pill_spacing = config["pill_spacing"]
    rect_width = config["rect_width"]
    rect_height = config["rect_height"]
    oval_size = config["oval_size"]
    overlap = config["overlap"]
    group_margin = config["group_margin"]
    top_margin = config["top_margin"]
    left_margin = config["left_margin"]
    bottom_margin = config["bottom_margin"]
    right_margin = config["right_margin"]
    boxes_per_row_first = config["boxes_per_row_first"]
    boxes_per_row_other = config["boxes_per_row_other"]

    # Group the Excel data by 'Category'
    groups = df.groupby("Category")
    
    # Base calculation for scaling
    forced_box_content_width = (1 * (rect_width + overlap)) + (1 * group_margin) + 2 * group_margin
    three_box_width = ((slide_width - forced_box_content_width) / 3) - 6 * group_margin

    # Build boxes list (one for each group)
    group_boxes = []
    for group_name, group_df in groups:
        rows = group_df.to_dict("records")
        n = len(rows)
        cols = min(2, n)  # Use up to 2 columns
        rows_needed = math.ceil(n / cols)

        content_width  = cols * (rect_width + overlap) + (cols - 1) * group_margin
        content_height = rows_needed * pill_spacing

        box_width  = three_box_width + 2 * group_margin
        box_height = content_height + 2 * group_margin

        group_boxes.append({
            "Category": group_name,
            "rows": rows,
            "n": n,
            "cols": cols,
            "box_width": box_width,
            "box_height": box_height
        })

    # Sort groups in descending order by number of pills ("n")
    group_boxes.sort(key=lambda gb: gb["n"], reverse=True)

    # Force the 4th box in the first row to have 1 column (if at least 4 groups exist)
    if len(group_boxes) >= 4:
        forced_box = group_boxes[3]  # 0-based index => 4th box
        forced_box["cols"] = 1
        n = forced_box["n"]
        content_width = (rect_width + overlap)  # Only one column; no inter-col margin
        content_height = n * pill_spacing
        forced_box["box_width"] = content_width + 2 * group_margin
        forced_box["box_height"] = content_height + 2 * group_margin

    # Arrange boxes into rows
    rows_list = []
    current_row = []
    row_count = 0
    for idx, gb in enumerate(group_boxes):
        current_row.append(gb)
        if row_count == 0:
            boxes_in_this_row = boxes_per_row_first
        else:
            boxes_in_this_row = boxes_per_row_other

        # When row is full or last box reached
        if (len(current_row) == boxes_in_this_row) or (idx == len(group_boxes) - 1):
            row_max_height = max(b["box_height"] for b in current_row)
            # Unify each box in row to have the same height
            for b in current_row:
                b["box_height"] = row_max_height
            rows_list.append(current_row)
            current_row = []
            row_count += 1

    # Horizontal scaling to equalize row widths
    row_metrics = []
    for row in rows_list:
        row_width = sum(b["box_width"] for b in row) + group_margin * (len(row) - 1)
        row_height = row[0]["box_height"]
        row_metrics.append({"boxes": row, "row_width": row_width, "row_height": row_height})

    # Calculate and apply scaling for each row
    max_row_width = max(r["row_width"] for r in row_metrics)
    for r in row_metrics:
        if r["row_width"] < max_row_width:
            scale = max_row_width / r["row_width"]
            for b in r["boxes"]:
                b["box_width"] *= scale
            r["row_width"] = max_row_width

    # Distribute extra vertical space across rows
    num_rows = len(row_metrics)
    total_used_height = sum(r["row_height"] for r in row_metrics) + group_margin * (num_rows - 1)
    remaining_height = slide_height - top_margin - bottom_margin - total_used_height
    extra_half = remaining_height / 2.0
    extra_per_row = extra_half / num_rows
    for r in row_metrics:
        r["row_height"] += extra_per_row
        for b in r["boxes"]:
            b["box_height"] = r["row_height"]

    # Final placement: compute (x,y) for each group box
    placements = []
    y_cursor = top_margin
    for r in row_metrics:
        row_width = r["row_width"]
        row_height = r["row_height"]
        x_cursor = left_margin + ((slide_width - left_margin - right_margin) - row_width) / 2
        for i, b in enumerate(r["boxes"]):
            b["x"] = x_cursor
            b["y"] = y_cursor
            b["box_index_in_row"] = i
            placements.append(b)
            x_cursor += b["box_width"] + group_margin
        y_cursor += row_height + group_margin

    # Draw the shapes
    text_padding = Inches(0.20)
    side_margin = Inches(0.005)

    for gb in placements:
        _draw_group_box(slide, gb, config, shortened_text_map, side_margin, text_padding)

def _draw_group_box(slide, gb, config, shortened_text_map, side_margin, text_padding):
    """Helper function to draw an individual group box"""
    new_x = gb["x"] + side_margin
    new_y = gb["y"] + side_margin
    new_w = gb["box_width"] - 2 * side_margin
    new_h = gb["box_height"] - 2 * side_margin

    # Draw the grey box
    grey_box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        new_x,
        new_y,
        new_w,
        new_h
    )
    grey_box.fill.solid()
    grey_box.fill.fore_color.rgb = RGBColor(220, 220, 220)
    grey_box.line.fill.solid()
    grey_box.line.fill.fore_color.rgb = RGBColor(0, 0, 0)
    grey_box.line.width = Pt(1)

    # Add the group label
    text_frame = grey_box.text_frame
    text_frame.clear()
    text_frame.vertical_anchor = MSO_ANCHOR.TOP
    p = text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = gb["Category"]
    run.font.size = Pt(10)
    run.font.bold = True
    run.font.name = "Arial"
    run.font.color.rgb = RGBColor(0, 0, 0)
    remove_bullets_from_paragraph(p)

    # Add banner value if applicable
    first_banner = next((row["Banner Value"] for row in gb["rows"] if row["Banner Value"] != "-"), None)
    format_info = analyze_format(first_banner) if first_banner else None
    banner_sum = sum(extract_number(row["Banner Value"]) for row in gb["rows"])

    if banner_sum > 0:
        _add_banner_value(
            slide, new_x, new_y, new_w, banner_sum, format_info,
            convert_to_millions=config.get("convert_to_millions", False),
            cost_decimal_places=config.get("cost_decimal_places", 1),
            cost_prefix=config.get("cost_prefix", "$"),
            cost_suffix=config.get("cost_suffix", "M")
        )

    # Calculate pill area
    cols_for_box = gb["cols"]
    pill_area_width = (cols_for_box * config["rect_width"]) + (cols_for_box - 1) * config["group_margin"]
    leftover = new_w - pill_area_width
    origin_x = new_x + (leftover / 2)
    origin_y = new_y + config["group_margin"] + text_padding

    # Draw each pill
    for i, row in enumerate(gb["rows"]):
        _draw_pill(slide, row, i, cols_for_box, origin_x, origin_y, config, shortened_text_map)

def _add_banner_value(slide, box_x, box_y, box_width, banner_sum, format_info, convert_to_millions=False, cost_decimal_places=1, cost_prefix="$", cost_suffix="M"):
    """Add banner value box to a group box"""
    # Format the banner value
    if convert_to_millions:
        banner_value = format_cost_to_millions(banner_sum, cost_decimal_places, cost_prefix, cost_suffix)
    else:
        banner_value = format_number(banner_sum, format_info)
    
    font_size = Pt(9)
    char_width = font_size * 0.6
    text_width = len(str(banner_value)) * char_width
    
    horizontal_padding = Inches(0.1)
    vertical_padding = Inches(0.05)
    banner_box_width = text_width + (2 * horizontal_padding)
    banner_box_height = font_size + (2 * vertical_padding)
    banner_box_right_margin = Inches(0.05)
    
    banner_box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        box_x + box_width - banner_box_width - banner_box_right_margin,
        box_y,
        banner_box_width,
        banner_box_height
    )
    banner_box.fill.solid()
    banner_box.fill.fore_color.rgb = RGBColor(240, 240, 240)
    banner_box.line.fill.solid()
    banner_box.line.fill.fore_color.rgb = RGBColor(128, 128, 128)
    banner_box.line.width = Pt(0.75)
    
    banner_text_frame = banner_box.text_frame
    banner_text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    banner_p = banner_text_frame.paragraphs[0]
    banner_p.alignment = PP_ALIGN.CENTER
    banner_run = banner_p.add_run()
    banner_run.text = str(banner_value)
    banner_run.font.size = font_size
    banner_run.font.bold = True
    banner_run.font.name = "Arial"
    banner_run.font.color.rgb = RGBColor(0, 0, 0)
    remove_bullets_from_paragraph(banner_p)

def _draw_pill(slide, row, index, cols, origin_x, origin_y, config, shortened_text_map):
    """Draw an individual pill shape with text and count"""
    col_idx = index % cols
    row_idx = index // cols
    if cols == 1:
        pill_left = origin_x + 1.5 * config["group_margin"]
    else:
        pill_left = origin_x + col_idx * (config["rect_width"] + config["group_margin"] + (1.5 * config["overlap"]))
    pill_top = origin_y + row_idx * config["pill_spacing"]

    # Get pill data
    sub_cat = row["Sub-Category"]
    count = row["Count"]
    color_name = row["Color"]
    pill_hex = config["color_map"].get(color_name, "000000")
    text_hex = "FFFFFF"
    display_count = str(count) if count and str(count).strip() != "-" else "-"
    if display_count == "-":
        pill_hex = "FFFFFF"
        text_hex = "000000"

    # Draw pill shape
    pill = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        pill_left,
        pill_top,
        config["rect_width"],
        config["rect_height"]
    )
    _set_pill_style(pill, pill_hex)

    # Add text to pill
    tf = pill.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    p = tf.paragraphs[0]
    _set_text_style(p, text_hex)

    # Handle text fitting and shortening
    available_width_pt = (config["rect_width"] - config["oval_size"] * 1.2) / Inches(1) * 72
    font_size_pt = 8
    _handle_pill_text(p, sub_cat, available_width_pt, font_size_pt, shortened_text_map)

    # Add count circle
    _add_count_circle(slide, pill_left, pill_top, config["oval_size"], config["overlap"], display_count)

def _set_pill_style(pill, hex_color):
    """Set the style for a pill shape"""
    pill.fill.solid()
    pill.fill.fore_color.rgb = RGBColor(
        int(hex_color[0:2], 16),
        int(hex_color[2:4], 16),
        int(hex_color[4:6], 16)
    )
    pill.line.fill.solid()
    pill.line.fill.fore_color.rgb = RGBColor(191, 191, 191)
    pill.line.width = Pt(0.25)

def _set_text_style(paragraph, hex_color):
    """Set the style for text in a pill"""
    paragraph.font.size = Pt(8)
    paragraph.font.bold = True
    paragraph.font.name = "Arial"
    paragraph.font.color.rgb = RGBColor(
        int(hex_color[0:2], 16),
        int(hex_color[2:4], 16),
        int(hex_color[4:6], 16)
    )
    paragraph.alignment = PP_ALIGN.CENTER
    remove_bullets_from_paragraph(paragraph)

def _handle_pill_text(paragraph, text, available_width, font_size, shortened_text_map):
    """Handle text fitting and shortening for pill text"""
    result, needs_shortening = split_text_for_pill(text, available_width, font_size)
    reference_number = None
    
    # Initialize the counter in the map if it doesn't exist
    if '_counter' not in shortened_text_map:
        shortened_text_map['_counter'] = 1
    
    if not needs_shortening and isinstance(result, tuple):
        first_part, second_part = result
        run = paragraph.add_run()
        run.text = first_part
        paragraph.add_line_break()
        run = paragraph.add_run()
        run.text = second_part
        return

    if needs_shortening:
        replacements = {
            "Software Development & Management": "Software Dev. & Management",
            "Cloud Services & SaaS Management": "Cloud Services & SaaS Mgmt.",
            "IT Service Management & Support": "IT Service Mgmt. & Support",
            "Performance, Training & Change Management": "Perf, Training & Change Mgmt.",
            "Security Education & Awareness": "Security Edu. & Awareness",
            "Internal Communication & Employee Engagement": "Internal Comm. & Employee Engmt.",
            "Customer Engagement & Digital Experience": "Cust. Engagement & Digital Experience",
            "Brand, Product & E-commerce Management": "Brand, Product & E-Comm Mgmt."
        }

        for original, short in replacements.items():
            if text == original:
                if original in shortened_text_map:
                    shortened_text = short
                    reference_number = shortened_text_map[original]
                else:
                    shortened_text = short
                    reference_number = f"[{shortened_text_map['_counter']}]"
                    shortened_text_map[original] = reference_number
                    shortened_text_map['_counter'] += 1
                break

        if 'shortened_text' in locals():
            result, needs_split = split_text_for_pill(shortened_text, available_width, font_size)
            if not needs_split and isinstance(result, tuple):
                first_part, second_part = result
                run = paragraph.add_run()
                run.text = first_part
                paragraph.add_line_break()
                run = paragraph.add_run()
                run.text = second_part
                text = None

    if text is not None:
        run = paragraph.add_run()
        run.text = text if 'shortened_text' not in locals() else shortened_text

    if reference_number:
        sup_run = paragraph.add_run()
        sup_run.text = reference_number
        sup_run.font.size = Pt(6)
        sup_run.font.bold = True
        sup_run.font.name = "Arial"
        sup_run.font.baseline = 30000
    
    # Ensure no bullets on this paragraph
    remove_bullets_from_paragraph(paragraph)

def _add_count_circle(slide, pill_left, pill_top, oval_size, overlap, display_count):
    """Add a count circle to a pill"""
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        pill_left - overlap,
        pill_top,
        oval_size,
        oval_size
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(0, 0, 0)
    circle.line.fill.background()

    textbox = slide.shapes.add_textbox(pill_left - overlap, pill_top, oval_size, oval_size)
    tf_box = textbox.text_frame
    tf_box.clear()
    tf_box.margin_left = 0
    tf_box.margin_right = 0
    p_box = tf_box.paragraphs[0]
    p_box.text = display_count
    p_box.font.size = Pt(10)
    p_box.font.bold = True
    p_box.font.name = "Arial"
    p_box.font.color.rgb = RGBColor(255, 255, 255)
    p_box.alignment = PP_ALIGN.CENTER
    tf_box.vertical_anchor = MSO_ANCHOR.MIDDLE
    remove_bullets_from_paragraph(p_box)
