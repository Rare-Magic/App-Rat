"""
Generate PowerPoint from output.csv. Called by Flask after Map CMDB.
"""
from pathlib import Path
import pandas as pd
from pptx import Presentation
from pptx.util import Inches

from .PPT_components import draw_group_boxes, add_slide_notes
from .PPT_components.presentation import remove_bullets_from_paragraph
from . import config as ppt_config

INPUT_CSV_PATH = ppt_config.INPUT_CSV_PATH
OUTPUT_PPT_PATH = ppt_config.OUTPUT_PPT_PATH
DEFAULT_THEME_PATH = ppt_config.DEFAULT_THEME_PATH
TEMPLATE_SLIDE_IDX = ppt_config.TEMPLATE_SLIDE_IDX
EXPECTED_COLUMNS = ppt_config.EXPECTED_COLUMNS
COLOR_RANGES = ppt_config.COLOR_RANGES
COLOR_MAP = ppt_config.COLOR_MAP
PILL_SPACING = ppt_config.PILL_SPACING
RECT_WIDTH = ppt_config.RECT_WIDTH
RECT_HEIGHT = ppt_config.RECT_HEIGHT
OVAL_SIZE = ppt_config.OVAL_SIZE
OVERLAP = ppt_config.OVERLAP
GROUP_MARGIN = ppt_config.GROUP_MARGIN
TOP_MARGIN = ppt_config.TOP_MARGIN
LEFT_MARGIN = ppt_config.LEFT_MARGIN
BOTTOM_MARGIN = ppt_config.BOTTOM_MARGIN
RIGHT_MARGIN = ppt_config.RIGHT_MARGIN
BOXES_PER_ROW_FIRST = ppt_config.BOXES_PER_ROW_FIRST
BOXES_PER_ROW_OTHER = ppt_config.BOXES_PER_ROW_OTHER
SLIDE_TITLE = ppt_config.SLIDE_TITLE
SOURCE_TEXT = ppt_config.SOURCE_TEXT
CONVERT_TO_MILLIONS = ppt_config.CONVERT_TO_MILLIONS
COST_DECIMAL_PLACES = ppt_config.COST_DECIMAL_PLACES
COST_PREFIX = ppt_config.COST_PREFIX
COST_SUFFIX = ppt_config.COST_SUFFIX


def load_csv_data(csv_path):
    if not Path(csv_path).exists():
        raise FileNotFoundError(f"CSV file not found: {csv_path}")
    df = pd.read_csv(csv_path)
    missing_cols = [col for col in EXPECTED_COLUMNS if col not in df.columns]
    if missing_cols:
        raise ValueError(f"CSV is missing required columns: {missing_cols}")
    df = df[EXPECTED_COLUMNS]
    df['Count'] = df['Count'].apply(
        lambda x: int(x) if x != '-' and str(x).strip() != '-' else '-'
    )
    return df


def assign_colors_to_data(df):
    colors = []
    for count in df['Count']:
        if count == '-':
            colors.append("White")
        else:
            assigned = None
            for color, (min_val, max_val, _) in COLOR_RANGES.items():
                if min_val <= count <= max_val:
                    assigned = color
                    break
            colors.append(assigned if assigned else "Red")
    df['Color'] = colors
    return df


def build_config_dict():
    return {
        "color_map": COLOR_MAP,
        "color_ranges": {
            color: (min_val, max_val)
            for color, (min_val, max_val, _) in COLOR_RANGES.items()
        },
        "output_path": str(OUTPUT_PPT_PATH),
        "default_theme_path": DEFAULT_THEME_PATH,
        "template_slide_idx": TEMPLATE_SLIDE_IDX,
        "pill_spacing": Inches(PILL_SPACING),
        "rect_width": Inches(RECT_WIDTH),
        "rect_height": Inches(RECT_HEIGHT),
        "oval_size": Inches(OVAL_SIZE),
        "overlap": Inches(OVERLAP),
        "group_margin": Inches(GROUP_MARGIN),
        "top_margin": Inches(TOP_MARGIN),
        "left_margin": Inches(LEFT_MARGIN),
        "bottom_margin": Inches(BOTTOM_MARGIN),
        "right_margin": Inches(RIGHT_MARGIN),
        "boxes_per_row_first": BOXES_PER_ROW_FIRST,
        "boxes_per_row_other": BOXES_PER_ROW_OTHER,
        "convert_to_millions": CONVERT_TO_MILLIONS,
        "cost_decimal_places": COST_DECIMAL_PLACES,
        "cost_prefix": COST_PREFIX,
        "cost_suffix": COST_SUFFIX
    }


def generate_slide_notes(df, config):
    notes = "Slide Summary:\n\n"
    total_sum = df.loc[df['Count'] != '-', 'Count'].astype(int).sum()
    notes += f"Total Sum across all Sub-Categories: {total_sum}\n\n"
    notes += "Sums by Category:\n"
    for cat in sorted(df['Category'].unique()):
        cat_sum = df[df['Category'] == cat].loc[df['Count'] != '-', 'Count'].astype(int).sum()
        notes += f"{cat}: {cat_sum}\n"
    notes += "\nColor Ranges:\n"
    for color, (min_val, max_val) in config["color_ranges"].items():
        if color == "Red":
            notes += f"- {color}: ≥ {min_val}\n"
        elif color != "White":
            notes += f"- {color}: {min_val} to {max_val}\n"
    return notes


def generate_presentation():
    """Generate output.pptx from output.csv. Returns True on success, False on error."""
    try:
        df = load_csv_data(INPUT_CSV_PATH)
        df = assign_colors_to_data(df)
        config = build_config_dict()
        if not Path(config["default_theme_path"]).exists():
            raise FileNotFoundError(f"Template not found: {config['default_theme_path']}")
        prs = Presentation(str(config["default_theme_path"]))

        def _find_template_source_shape(presentation_obj):
            def _shape_has_source_text(shape):
                try:
                    if hasattr(shape, 'text') and shape.text:
                        if 'source' in shape.text.lower() or 'note' in shape.text.lower():
                            return True
                except Exception:
                    pass
                return False
            for layout in presentation_obj.slide_layouts:
                try:
                    for shape in layout.shapes:
                        if _shape_has_source_text(shape):
                            return shape
                except Exception:
                    continue
            for s in presentation_obj.slides:
                try:
                    for shape in s.shapes:
                        if _shape_has_source_text(shape):
                            return shape
                except Exception:
                    continue
            for s in presentation_obj.slides:
                try:
                    ns_text = s.notes_slide.notes_text_frame.text
                    if ns_text and ("source" in ns_text.lower() or "note" in ns_text.lower()):
                        return {'text': ns_text}
                except Exception:
                    continue
            return None

        template_shape = _find_template_source_shape(prs)
        while len(prs.slides) > 0:
            rId = prs.slides._sldIdLst[0].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[0]
        slide = prs.slides.add_slide(prs.slide_layouts[config["template_slide_idx"]])
        shortened_text_map = {}
        draw_group_boxes(prs, slide, df, config, shortened_text_map, slide_title=SLIDE_TITLE)
        notes = generate_slide_notes(df, config)
        if shortened_text_map:
            notes += "\nText Shortenings:\n"
            reference_items = {k: v for k, v in shortened_text_map.items() if k != '_counter'}
            sorted_items = sorted(reference_items.items(), key=lambda x: int(x[1].strip('[]')))
            for original, ref in sorted_items:
                notes += f"{ref}: {original}; "
        if template_shape and SOURCE_TEXT:
            notes += "\n\n" + SOURCE_TEXT
            if not isinstance(template_shape, dict):
                try:
                    src_shape = template_shape
                    left = getattr(src_shape, 'left', None)
                    top = getattr(src_shape, 'top', None)
                    width = getattr(src_shape, 'width', None)
                    height = getattr(src_shape, 'height', None)
                    if None not in (left, top, width, height) and hasattr(src_shape, 'text_frame'):
                        src_tf = src_shape.text_frame
                        tb = slide.shapes.add_textbox(left, top, width, height)
                        new_tf = tb.text_frame
                        try:
                            new_tf.margin_left = src_tf.margin_left
                            new_tf.margin_right = src_tf.margin_right
                            new_tf.vertical_anchor = src_tf.vertical_anchor
                        except Exception:
                            pass
                        try:
                            new_tf.clear()
                        except Exception:
                            new_tf.text = ''
                        template_source_para = None
                        for src_p in src_tf.paragraphs:
                            if "source:" in src_p.text.lower():
                                template_source_para = src_p
                                break
                        p = new_tf.paragraphs[0] if len(new_tf.paragraphs) > 0 and new_tf.paragraphs[0].text == '' else new_tf.add_paragraph()
                        p.text = "Source: " + SOURCE_TEXT
                        # Make source text small (Arial 8)
                        try:
                            for run in p.runs:
                                run.font.size = Pt(8)
                                run.font.name = ppt_config.FONT_NAME
                        except Exception:
                            pass
                        if template_source_para:
                            try:
                                p.alignment = template_source_para.alignment
                                remove_bullets_from_paragraph(p)
                            except Exception:
                                pass
                        else:
                            try:
                                remove_bullets_from_paragraph(p)
                            except Exception:
                                pass
                except Exception:
                    pass
        add_slide_notes(slide, notes)

        # Ensure all “Source:” text is small (Arial 8) even if it comes from the template
        try:
            from pptx.util import Pt as _Pt
            for slide_obj in prs.slides:
                # Shapes on slide
                for shape in slide_obj.shapes:
                    if not hasattr(shape, "text_frame"):
                        continue
                    tf = shape.text_frame
                    for para in tf.paragraphs:
                        if para.text and "source:" in para.text.lower():
                            for run in para.runs:
                                run.font.size = _Pt(8)
                                run.font.name = ppt_config.FONT_NAME
                # Notes
                if hasattr(slide_obj, "notes_slide") and slide_obj.notes_slide:
                    ns_tf = slide_obj.notes_slide.notes_text_frame
                    for para in ns_tf.paragraphs:
                        if para.text and "source:" in para.text.lower():
                            for run in para.runs:
                                run.font.size = _Pt(8)
                                run.font.name = ppt_config.FONT_NAME
        except Exception:
            pass

        OUTPUT_PPT_PATH.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(OUTPUT_PPT_PATH))
        return True
    except Exception:
        return False
